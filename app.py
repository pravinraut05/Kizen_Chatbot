import streamlit as st
import pandas as pd
import PyPDF2
import docx
from pptx import Presentation
import io
import re
from sentence_transformers import SentenceTransformer
import numpy as np
from sklearn.metrics.pairwise import cosine_similarity
from typing import List, Dict, Any

# Configure page
st.set_page_config(
    page_title="Document Q&A Chatbot",
    page_icon="🤖",
    layout="wide"
)

# Initialize session state
if 'processed_content' not in st.session_state:
    st.session_state.processed_content = ""
if 'embeddings' not in st.session_state:
    st.session_state.embeddings = None
if 'chunks' not in st.session_state:
    st.session_state.chunks = []
if 'chat_history' not in st.session_state:
    st.session_state.chat_history = []

class DocumentProcessor:
    def __init__(self):
        self.model = SentenceTransformer('all-MiniLM-L6-v2')
    
    def extract_text_from_pdf(self, file) -> str:
        """Extract text from PDF file"""
        try:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text
        except Exception as e:
            st.error(f"Error reading PDF: {str(e)}")
            return ""
    
    def extract_text_from_docx(self, file) -> str:
        """Extract text from Word document"""
        try:
            doc = docx.Document(file)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            st.error(f"Error reading Word document: {str(e)}")
            return ""
    
    def extract_text_from_pptx(self, file) -> str:
        """Extract text from PowerPoint presentation"""
        try:
            prs = Presentation(file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            st.error(f"Error reading PowerPoint: {str(e)}")
            return ""
    
    def extract_text_from_csv(self, file) -> str:
        """Extract text from CSV file"""
        try:
            df = pd.read_csv(file)
            # Convert DataFrame to readable text format
            text = "CSV Data Summary:\n"
            text += f"Columns: {', '.join(df.columns)}\n"
            text += f"Number of rows: {len(df)}\n\n"
            text += "Data Preview:\n"
            text += df.head(10).to_string() + "\n\n"
            
            # Add column statistics for numeric columns
            numeric_cols = df.select_dtypes(include=[np.number]).columns
            if len(numeric_cols) > 0:
                text += "Numeric Column Statistics:\n"
                text += df[numeric_cols].describe().to_string() + "\n"
            
            return text
        except Exception as e:
            st.error(f"Error reading CSV: {str(e)}")
            return ""
    
    def extract_text_from_excel(self, file) -> str:
        """Extract text from Excel file"""
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file)
            text = f"Excel file with {len(excel_file.sheet_names)} sheet(s):\n"
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file, sheet_name=sheet_name)
                text += f"\nSheet: {sheet_name}\n"
                text += f"Columns: {', '.join(df.columns)}\n"
                text += f"Number of rows: {len(df)}\n"
                text += "Data Preview:\n"
                text += df.head(5).to_string() + "\n"
            
            return text
        except Exception as e:
            st.error(f"Error reading Excel: {str(e)}")
            return ""
    
    def process_file(self, file) -> str:
        """Process uploaded file based on its type"""
        file_type = file.type
        
        if file_type == "application/pdf":
            return self.extract_text_from_pdf(file)
        elif file_type in ["application/vnd.openxmlformats-officedocument.wordprocessingml.document", 
                          "application/msword"]:
            return self.extract_text_from_docx(file)
        elif file_type in ["application/vnd.openxmlformats-officedocument.presentationml.presentation",
                          "application/vnd.ms-powerpoint"]:
            return self.extract_text_from_pptx(file)
        elif file_type == "text/csv":
            return self.extract_text_from_csv(file)
        elif file_type in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                          "application/vnd.ms-excel"]:
            return self.extract_text_from_excel(file)
        else:
            st.error(f"Unsupported file type: {file_type}")
            return ""
    
    def chunk_text(self, text: str, chunk_size: int = 500) -> List[str]:
        """Split text into chunks for better processing"""
        words = text.split()
        chunks = []
        current_chunk = []
        
        for word in words:
            current_chunk.append(word)
            if len(' '.join(current_chunk)) > chunk_size:
                chunks.append(' '.join(current_chunk[:-1]))
                current_chunk = [word]
        
        if current_chunk:
            chunks.append(' '.join(current_chunk))
        
        return chunks
    
    def create_embeddings(self, chunks: List[str]) -> np.ndarray:
        """Create embeddings for text chunks"""
        try:
            embeddings = self.model.encode(chunks)
            return embeddings
        except Exception as e:
            st.error(f"Error creating embeddings: {str(e)}")
            return np.array([])
    
    def find_relevant_chunks(self, query: str, chunks: List[str], embeddings: np.ndarray, top_k: int = 3) -> List[str]:
        """Find most relevant chunks for the query"""
        if len(embeddings) == 0:
            return chunks[:top_k] if chunks else []
        
        query_embedding = self.model.encode([query])
        similarities = cosine_similarity(query_embedding, embeddings)[0]
        
        top_indices = np.argsort(similarities)[-top_k:][::-1]
        return [chunks[i] for i in top_indices if i < len(chunks)]

class SimpleChatbot:
    """Simple rule-based chatbot for basic Q&A"""
    
    def __init__(self, content: str):
        self.content = content.lower()
        self.original_content = content
    
    def answer_question(self, question: str) -> str:
        """Generate answer based on content analysis"""
        question = question.lower()
        
        # Simple keyword matching and response generation
        if any(word in question for word in ['how many', 'count', 'total', 'number']):
            return self._handle_count_question(question)
        elif any(word in question for word in ['what', 'define', 'explain']):
            return self._handle_what_question(question)
        elif any(word in question for word in ['when', 'date', 'time']):
            return self._handle_when_question(question)
        elif any(word in question for word in ['where', 'location']):
            return self._handle_where_question(question)
        elif any(word in question for word in ['why', 'reason']):
            return self._handle_why_question(question)
        elif any(word in question for word in ['summary', 'summarize', 'overview']):
            return self._provide_summary()
        else:
            return self._general_search(question)
    
    def _handle_count_question(self, question: str) -> str:
        """Handle counting questions"""
        lines = self.original_content.split('\n')
        words = self.original_content.split()
        
        if 'rows' in question or 'records' in question:
            row_info = [line for line in lines if 'rows' in line.lower()]
            if row_info:
                return f"Based on the document: {row_info[0]}"
        
        return f"The document contains approximately {len(words)} words and {len(lines)} lines of content."
    
    def _handle_what_question(self, question: str) -> str:
        """Handle 'what' questions"""
        # Extract key terms from question
        question_words = question.split()
        key_terms = [word for word in question_words if len(word) > 3 and word not in ['what', 'does', 'mean', 'this']]
        
        relevant_sentences = []
        for sentence in self.original_content.split('.'):
            if any(term in sentence.lower() for term in key_terms):
                relevant_sentences.append(sentence.strip())
        
        if relevant_sentences:
            return f"Based on the document: {'. '.join(relevant_sentences[:2])}."
        else:
            return "I couldn't find specific information about that in the document. Could you try rephrasing your question?"
    
    def _handle_when_question(self, question: str) -> str:
        """Handle 'when' questions"""
        # Look for dates and time references
        date_patterns = [r'\d{4}', r'\d{1,2}/\d{1,2}/\d{4}', r'\d{1,2}-\d{1,2}-\d{4}']
        dates_found = []
        
        for pattern in date_patterns:
            dates_found.extend(re.findall(pattern, self.original_content))
        
        if dates_found:
            return f"The document mentions these dates/years: {', '.join(set(dates_found))}"
        else:
            return "I couldn't find specific date information in the document."
    
    def _handle_where_question(self, question: str) -> str:
        """Handle 'where' questions"""
        # Look for location indicators
        location_keywords = ['country', 'city', 'state', 'location', 'address', 'region']
        relevant_info = []
        
        for sentence in self.original_content.split('.'):
            if any(keyword in sentence.lower() for keyword in location_keywords):
                relevant_info.append(sentence.strip())
        
        if relevant_info:
            return f"Location information from the document: {'. '.join(relevant_info[:2])}"
        else:
            return "I couldn't find specific location information in the document."
    
    def _handle_why_question(self, question: str) -> str:
        """Handle 'why' questions"""
        reason_keywords = ['because', 'due to', 'reason', 'cause', 'since', 'as a result']
        relevant_info = []
        
        for sentence in self.original_content.split('.'):
            if any(keyword in sentence.lower() for keyword in reason_keywords):
                relevant_info.append(sentence.strip())
        
        if relevant_info:
            return f"Reasoning from the document: {'. '.join(relevant_info[:2])}"
        else:
            return "I couldn't find specific reasoning or explanations in the document."
    
    def _provide_summary(self) -> str:
        """Provide a summary of the document"""
        lines = self.original_content.split('\n')
        non_empty_lines = [line.strip() for line in lines if line.strip()]
        
        if len(non_empty_lines) > 10:
            summary_lines = non_empty_lines[:5] + ['...'] + non_empty_lines[-2:]
        else:
            summary_lines = non_empty_lines
        
        return f"Document Summary:\n" + '\n'.join(summary_lines[:200])  # Limit length
    
    def _general_search(self, question: str) -> str:
        """General search through content"""
        question_words = [word for word in question.split() if len(word) > 3]
        relevant_sentences = []
        
        for sentence in self.original_content.split('.'):
            if any(word in sentence.lower() for word in question_words):
                relevant_sentences.append(sentence.strip())
        
        if relevant_sentences:
            return f"Relevant information: {'. '.join(relevant_sentences[:3])}"
        else:
            return "I couldn't find information directly related to your question. Try asking about specific topics mentioned in the document."

def main():
    st.title("🤖 Document Q&A Chatbot")
    st.markdown("Upload documents (PDF, Word, PowerPoint, CSV, Excel) and ask questions about their content!")
    
    # Initialize processor
    processor = DocumentProcessor()
    
    # Sidebar for file upload
    with st.sidebar:
        st.header("📁 Upload Documents")
        uploaded_files = st.file_uploader(
            "Choose files",
            type=['pdf', 'docx', 'doc', 'pptx', 'ppt', 'csv', 'xlsx', 'xls'],
            accept_multiple_files=True
        )
        
        if uploaded_files:
            st.success(f"Uploaded {len(uploaded_files)} file(s)")
            
            if st.button("Process Files"):
                with st.spinner("Processing files..."):
                    all_content = ""
                    for file in uploaded_files:
                        st.write(f"Processing: {file.name}")
                        content = processor.process_file(file)
                        all_content += f"\n\n--- Content from {file.name} ---\n{content}"
                    
                    if all_content.strip():
                        st.session_state.processed_content = all_content
                        st.session_state.chunks = processor.chunk_text(all_content)
                        st.session_state.embeddings = processor.create_embeddings(st.session_state.chunks)
                        st.success("Files processed successfully!")
                    else:
                        st.error("No content could be extracted from the files.")
    
    # Main chat interface
    col1, col2 = st.columns([2, 1])
    
    with col1:
        st.header("💬 Chat with your documents")
        
        # Display chat history
        for i, (question, answer) in enumerate(st.session_state.chat_history):
            with st.container():
                st.write(f"**You:** {question}")
                st.write(f"**Bot:** {answer}")
                st.divider()
        
        # Chat input
        if st.session_state.processed_content:
            question = st.text_input("Ask a question about your documents:", key="question_input")
            
            if st.button("Send") and question:
                with st.spinner("Thinking..."):
                    # Initialize chatbot with processed content
                    chatbot = SimpleChatbot(st.session_state.processed_content)
                    answer = chatbot.answer_question(question)
                    
                    # Add to chat history
                    st.session_state.chat_history.append((question, answer))
                    st.rerun()
        else:
            st.info("Please upload and process documents first using the sidebar.")
    
    with col2:
        st.header("📊 Document Info")
        if st.session_state.processed_content:
            content_length = len(st.session_state.processed_content)
            word_count = len(st.session_state.processed_content.split())
            chunk_count = len(st.session_state.chunks)
            
            st.metric("Content Length", f"{content_length:,} characters")
            st.metric("Word Count", f"{word_count:,} words")
            st.metric("Text Chunks", chunk_count)
            
            # Clear chat button
            if st.button("Clear Chat History"):
                st.session_state.chat_history = []
                st.rerun()
        else:
            st.info("No documents processed yet.")

if __name__ == "__main__":
    main()
